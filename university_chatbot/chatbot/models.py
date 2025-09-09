from django.db import models
from django.contrib.auth.models import User
from django.utils.translation import gettext_lazy as _
import os
import PyPDF2
from pptx import Presentation
import io

class Document(models.Model):
    SEMESTER_CHOICES = [
        (1, '1st Semester'),
        (2, '2nd Semester'),
        (3, '3rd Semester'),
        (4, '4th Semester'),
        (5, '5th Semester'),
        (6, '6th Semester'),
        (7, '7th Semester'),
        (8, '8th Semester'),
    ]
    
    DOC_TYPE_CHOICES = [
        ('notes', 'Notes'),
        ('ppt', 'PPT'),
        ('syllabus', 'Syllabus'),
        ('circular', 'Circular'),
        ('assignment', 'Assignment'),
        ('question_paper', 'Question Paper'),
    ]
    
    title = models.CharField(max_length=200)
    description = models.TextField(blank=True)
    file = models.FileField(upload_to='documents/%Y/%m/%d/')
    semester = models.IntegerField(choices=SEMESTER_CHOICES)
    subject = models.CharField(max_length=100)
    unit = models.IntegerField(blank=True, null=True)
    doc_type = models.CharField(max_length=50, choices=DOC_TYPE_CHOICES)
    uploaded_by = models.ForeignKey(User, on_delete=models.SET_NULL, null=True, blank=True)
    uploaded_at = models.DateTimeField(auto_now_add=True)
    is_active = models.BooleanField(default=True)
    extracted_text = models.TextField(blank=True)
    
    def __str__(self):
        return self.title
    
    def filename(self):
        return os.path.basename(self.file.name)
    
    def extension(self):
        name, extension = os.path.splitext(self.file.name)
        return extension.lower()
    
    def save(self, *args, **kwargs):
        # Extract text when saving
        if self.file and not self.extracted_text:
            self.extracted_text = self.extract_text_from_file()
        super().save(*args, **kwargs)
    
    def extract_text_from_file(self):
        """
        Extract text content from uploaded files
        """
        try:
            file_extension = os.path.splitext(self.file.name)[1].lower()
            
            if file_extension == '.pdf':
                return self.extract_text_from_pdf()
            elif file_extension in ['.ppt', '.pptx']:
                return self.extract_text_from_ppt()
            else:
                return ""
                
        except Exception as e:
            print(f"Error extracting text: {e}")
            return ""
    
    def extract_text_from_pdf(self):
        """Extract text from PDF files"""
        text = ""
        try:
            # Save file to temporary location for reading
            with open(self.file.path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error reading PDF: {e}")
        return text
    
    def extract_text_from_ppt(self):
        """Extract text from PowerPoint files"""
        text = ""
        try:
            presentation = Presentation(self.file.path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error reading PPT: {e}")
        return text
    
    def delete(self, *args, **kwargs):
        # Delete the file from storage when the document is deleted
        self.file.delete()
        super().delete(*args, **kwargs)

class FAQ(models.Model):
    question = models.CharField(max_length=300)
    answer = models.TextField()
    category = models.CharField(max_length=100, choices=[
        ('admission', 'Admission'),
        ('examination', 'Examination'),
        ('scholarship', 'Scholarship'),
        ('hostel', 'Hostel'),
        ('general', 'General'),
    ])
    
    def __str__(self):
        return self.question

class Student(models.Model):
    user = models.OneToOneField(User, on_delete=models.CASCADE)
    enrollment_no = models.CharField(max_length=20, unique=True)
    semester = models.IntegerField()
    department = models.CharField(max_length=100)
    current_attendance = models.FloatField(default=0.0)
    
    def __str__(self):
        return f"{self.user.first_name} {self.user.last_name} ({self.enrollment_no})"

class Lecture(models.Model):
    DAY_CHOICES = [
        ('monday', 'Monday'),
        ('tuesday', 'Tuesday'),
        ('wednesday', 'Wednesday'),
        ('thursday', 'Thursday'),
        ('friday', 'Friday'),
        ('saturday', 'Saturday'),
    ]
    
    student = models.ForeignKey(Student, on_delete=models.CASCADE)
    day = models.CharField(max_length=10, choices=DAY_CHOICES)
    time_slot = models.CharField(max_length=50)
    subject = models.CharField(max_length=100)
    faculty = models.CharField(max_length=100)
    email = models.EmailField(blank=True)
    classroom = models.CharField(max_length=100)
    is_attended = models.BooleanField(default=False)
    date = models.DateField(null=True, blank=True)
    
    def __str__(self):
        return f"{self.subject} - {self.day} {self.time_slot}"

class AttendanceRecord(models.Model):
    student = models.ForeignKey(Student, on_delete=models.CASCADE)
    subject = models.CharField(max_length=100)
    total_classes = models.IntegerField(default=0)
    attended_classes = models.IntegerField(default=0)
    last_updated = models.DateField(auto_now=True)
    
    @property
    def percentage(self):
        if self.total_classes == 0:
            return 0
        return round((self.attended_classes / self.total_classes) * 100, 2)
    
    def __str__(self):
        return f"{self.student} - {self.subject}: {self.percentage}%"

class Timetable(models.Model):
    DAY_CHOICES = [
        ('monday', 'Monday'),
        ('tuesday', 'Tuesday'),
        ('wednesday', 'Wednesday'),
        ('thursday', 'Thursday'),
        ('friday', 'Friday'),
        ('saturday', 'Saturday'),
    ]
    
    semester = models.IntegerField()
    day = models.CharField(max_length=10, choices=DAY_CHOICES)
    time_slot = models.CharField(max_length=50)
    subject = models.CharField(max_length=100)
    faculty = models.CharField(max_length=100, blank=True)
    location = models.CharField(max_length=100, blank=True)